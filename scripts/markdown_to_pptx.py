#!/usr/bin/env python3
"""
Markdown to PowerPoint Converter

This script converts markdown files into professionally formatted PowerPoint presentations.
It supports various slide types, custom styling, and advanced formatting options.
"""

import os
import re
import sys
import logging
from typing import List, Optional, Dict, Any
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import markdown
from bs4 import BeautifulSoup
import yaml

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

@dataclass
class SlideTheme:
    """Theme configuration for presentation slides."""
    background_color: RGBColor = RGBColor(255, 255, 255)
    title_color: RGBColor = RGBColor(0, 32, 96)
    text_color: RGBColor = RGBColor(0, 0, 0)
    accent_color: RGBColor = RGBColor(0, 112, 192)
    title_font: str = 'Calibri'
    body_font: str = 'Calibri'
    title_size: int = 44
    subtitle_size: int = 32
    body_size: int = 18
    bullet_size: int = 18

class PresentationBuilder:
    """Builds PowerPoint presentations from markdown content."""
    
    def __init__(self, theme: Optional[SlideTheme] = None):
        """Initialize the presentation builder with optional theme."""
        self.prs = Presentation()
        self.theme = theme or SlideTheme()
        
        # Slide layouts
        self.title_slide_layout = self.prs.slide_layouts[0]
        self.content_slide_layout = self.prs.slide_layouts[1]
        self.section_slide_layout = self.prs.slide_layouts[2]
        self.image_slide_layout = self.prs.slide_layouts[5]
        self.two_content_layout = self.prs.slide_layouts[3]
        
        self.current_slide = None
        self._setup_slide_master()

    def _setup_slide_master(self) -> None:
        """Configure the slide master with theme settings."""
        slide_master = self.prs.slide_master
        
        # Set background
        background = slide_master.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.theme.background_color
        
        # Set theme colors
        for shape in slide_master.shapes:
            if hasattr(shape, 'text'):
                for paragraph in shape.text_frame.paragraphs:
                    paragraph.font.name = self.theme.body_font
                    paragraph.font.color.rgb = self.theme.text_color

    def create_title_slide(self, title: str, subtitle: str = "") -> None:
        """Create a title slide with optional subtitle."""
        try:
            slide = self.prs.slides.add_slide(self.title_slide_layout)
            title_shape = slide.shapes.title
            subtitle_shape = slide.placeholders[1]

            # Style title
            title_shape.text = title
            title_frame = title_shape.text_frame
            title_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            paragraph = title_frame.paragraphs[0]
            paragraph.font.size = Pt(self.theme.title_size)
            paragraph.font.bold = True
            paragraph.font.name = self.theme.title_font
            paragraph.font.color.rgb = self.theme.title_color
            paragraph.alignment = PP_ALIGN.CENTER

            # Style subtitle
            subtitle_shape.text = subtitle
            subtitle_frame = subtitle_shape.text_frame
            subtitle_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            paragraph = subtitle_frame.paragraphs[0]
            paragraph.font.size = Pt(self.theme.subtitle_size)
            paragraph.font.name = self.theme.body_font
            paragraph.font.color.rgb = self.theme.text_color
            paragraph.alignment = PP_ALIGN.CENTER

        except Exception as e:
            logger.error(f"Error creating title slide: {str(e)}")
            raise

    def create_section_slide(self, title: str) -> None:
        """Create a section divider slide."""
        try:
            slide = self.prs.slides.add_slide(self.section_slide_layout)
            title_shape = slide.shapes.title
            title_shape.text = title

            # Style section title
            frame = title_shape.text_frame
            frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            paragraph = frame.paragraphs[0]
            paragraph.font.size = Pt(self.theme.title_size)
            paragraph.font.bold = True
            paragraph.font.name = self.theme.title_font
            paragraph.font.color.rgb = self.theme.title_color
            paragraph.alignment = PP_ALIGN.CENTER

            # Add decorative line
            line = slide.shapes.add_shape(
                MSO_SHAPE.LINE_INVERSE,
                Inches(2), Inches(4),
                Inches(6), Inches(0)
            )
            line.line.color.rgb = self.theme.accent_color
            line.line.width = Pt(3)

        except Exception as e:
            logger.error(f"Error creating section slide: {str(e)}")
            raise

    def create_content_slide(self, title: str, content: List[str]) -> None:
        """Create a content slide with bullet points."""
        try:
            slide = self.prs.slides.add_slide(self.content_slide_layout)
            title_shape = slide.shapes.title
            content_shape = slide.placeholders[1]

            # Style title
            title_shape.text = title
            title_frame = title_shape.text_frame
            paragraph = title_frame.paragraphs[0]
            paragraph.font.size = Pt(self.theme.subtitle_size)
            paragraph.font.bold = True
            paragraph.font.name = self.theme.title_font
            paragraph.font.color.rgb = self.theme.title_color

            # Add content
            text_frame = content_shape.text_frame
            text_frame.clear()  # Clear default content
            
            for line in content:
                p = text_frame.add_paragraph()
                p.text = line.strip()
                p.font.size = Pt(self.theme.body_size)
                p.font.name = self.theme.body_font
                p.font.color.rgb = self.theme.text_color
                
                # Handle bullet points
                if line.strip().startswith(('•', '-', '*')):
                    p.text = line.strip().lstrip('•-* ')
                    p.level = 1
                    p.font.size = Pt(self.theme.bullet_size)

        except Exception as e:
            logger.error(f"Error creating content slide: {str(e)}")
            raise

    def create_image_slide(self, title: str, image_path: str, caption: str = "") -> None:
        """Create an image slide with optional caption."""
        try:
            slide = self.prs.slides.add_slide(self.image_slide_layout)
            title_shape = slide.shapes.title
            title_shape.text = title

            # Style title
            title_frame = title_shape.text_frame
            paragraph = title_frame.paragraphs[0]
            paragraph.font.size = Pt(self.theme.subtitle_size)
            paragraph.font.bold = True
            paragraph.font.name = self.theme.title_font
            paragraph.font.color.rgb = self.theme.title_color

            # Add image
            if os.path.exists(image_path):
                # Calculate image dimensions while maintaining aspect ratio
                img_width = Inches(8)
                left = Inches(1)
                top = Inches(2)
                
                slide.shapes.add_picture(
                    image_path,
                    left, top,
                    width=img_width
                )

                # Add caption if provided
                if caption:
                    left = Inches(1)
                    top = Inches(6.5)
                    width = Inches(8)
                    height = Inches(0.5)
                    
                    txBox = slide.shapes.add_textbox(left, top, width, height)
                    tf = txBox.text_frame
                    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
                    
                    p = tf.add_paragraph()
                    p.text = caption
                    p.alignment = PP_ALIGN.CENTER
                    p.font.size = Pt(14)
                    p.font.italic = True
                    p.font.name = self.theme.body_font
                    p.font.color.rgb = self.theme.text_color
            else:
                logger.warning(f"Image not found: {image_path}")

        except Exception as e:
            logger.error(f"Error creating image slide: {str(e)}")
            raise

def parse_markdown_metadata(content: str) -> tuple[Dict[str, Any], str]:
    """Parse YAML metadata from markdown content."""
    if content.startswith('---'):
        try:
            _, fm, content = content.split('---', 2)
            metadata = yaml.safe_load(fm)
            return metadata, content
        except:
            return {}, content
    return {}, content

def parse_theme_metadata(metadata: Dict[str, Any]) -> Optional[SlideTheme]:
    """Parse theme metadata and create SlideTheme object."""
    if not metadata.get('theme'):
        return None
        
    theme_data = metadata['theme']
    
    # Convert color arrays to RGBColor objects
    if 'title_color' in theme_data:
        theme_data['title_color'] = RGBColor(*theme_data['title_color'])
    if 'accent_color' in theme_data:
        theme_data['accent_color'] = RGBColor(*theme_data['accent_color'])
    if 'text_color' in theme_data:
        theme_data['text_color'] = RGBColor(*theme_data['text_color'])
    if 'background_color' in theme_data:
        theme_data['background_color'] = RGBColor(*theme_data['background_color'])
    
    return SlideTheme(**theme_data)

def convert_markdown_to_pptx(markdown_file: str, output_file: str, theme: Optional[SlideTheme] = None) -> None:
    """Convert markdown file to PowerPoint presentation."""
    try:
        logger.info(f"Converting {markdown_file} to PowerPoint...")
        
        # Read markdown content
        with open(markdown_file, 'r', encoding='utf-8') as f:
            content = f.read()

        # Parse metadata
        metadata, content = parse_markdown_metadata(content)
        
        # Create theme from metadata if provided
        if metadata.get('theme'):
            theme = parse_theme_metadata(metadata)

        # Split content into slides
        slides = content.split('---')

        # Initialize presentation builder
        builder = PresentationBuilder(theme)

        # Create title slide if metadata exists
        if metadata.get('title'):
            builder.create_title_slide(
                metadata['title'],
                metadata.get('subtitle', '')
            )

        # Process each slide
        for slide in slides:
            if not slide.strip():
                continue

            # Parse the slide content
            lines = slide.strip().split('\n')
            title = ""
            content = []
            image_path = None
            
            for line in lines:
                if line.startswith('# '):
                    title = line[2:].strip()
                elif line.startswith('## '):
                    if not title:
                        title = line[3:].strip()
                    else:
                        content.append(line[3:].strip())
                elif line.startswith('!['):
                    # Extract image path and caption from markdown image syntax
                    match = re.search(r'\!\[(.*?)\]\((.*?)\)', line)
                    if match:
                        caption = match.group(1)
                        image_path = match.group(2)
                        image_path = os.path.join(
                            os.path.dirname(markdown_file),
                            '..',
                            image_path
                        )
                elif line.strip():
                    content.append(line.strip())

            # Create appropriate slide type
            if title and not content and not image_path:
                builder.create_section_slide(title)
            elif image_path:
                builder.create_image_slide(title, image_path)
            elif title and content:
                builder.create_content_slide(title, content)
            elif title:
                builder.create_title_slide(title)

        # Save the presentation
        os.makedirs(os.path.dirname(output_file), exist_ok=True)
        builder.prs.save(output_file)
        logger.info(f"Presentation saved successfully: {output_file}")

    except Exception as e:
        logger.error(f"Error converting markdown to PowerPoint: {str(e)}")
        raise

def main():
    """Main entry point for the script."""
    try:
        # Get command line arguments
        if len(sys.argv) > 2:
            markdown_file = sys.argv[1]
            output_file = sys.argv[2]
        else:
            # Default paths
            base_dir = Path(__file__).parent.parent
            markdown_file = base_dir / 'docs' / 'test_presentation.md'
            output_file = base_dir / 'presentation' / 'test_presentation.pptx'
        
        # Create custom theme
        theme = SlideTheme(
            title_color=RGBColor(0, 32, 96),
            accent_color=RGBColor(0, 112, 192),
            title_font='Calibri',
            body_font='Calibri Light'
        )
        
        # Convert markdown to PowerPoint
        convert_markdown_to_pptx(str(markdown_file), str(output_file), theme)

    except Exception as e:
        logger.error(f"Script execution failed: {str(e)}")
        sys.exit(1)

if __name__ == "__main__":
    main()
