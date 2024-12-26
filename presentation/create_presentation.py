from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_algo360fx_presentation():
    prs = Presentation()
    
    # Title Slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]
    title.text = "Algo360FX"
    subtitle.text = "Advanced Forex Trading Analytics Platform"
    
    # Overview Slide
    overview_slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = overview_slide.shapes.title
    content = overview_slide.placeholders[1]
    title.text = "Platform Overview"
    content.text = """
    • Comprehensive Forex Trading Analysis Platform
    • Real-time Technical Analysis
    • Advanced Risk Management
    • Market Sentiment Analysis
    • News Integration
    • Automated Trading Signals
    """
    
    # Technical Analysis Slide
    tech_slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = tech_slide.shapes.title
    content = tech_slide.placeholders[1]
    title.text = "Technical Analysis Features"
    content.text = """
    • Multiple Technical Indicators:
        - SMA, EMA
        - RSI
        - MACD
        - Bollinger Bands
        - ATR
        - Stochastic
    • Custom Indicator Combinations
    • Trend Analysis
    """
    
    # Risk Management Slide
    risk_slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = risk_slide.shapes.title
    content = risk_slide.placeholders[1]
    title.text = "Risk Management System"
    content.text = """
    • Position Size Calculator
    • Risk-Reward Analysis
    • Maximum Drawdown Control
    • Stop Loss Optimization
    • Account Balance Protection
    • Risk Percentage per Trade
    • Position Risk Assessment
    """
    
    # Market Intelligence Slide
    market_slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = market_slide.shapes.title
    content = market_slide.placeholders[1]
    title.text = "Market Intelligence"
    content.text = """
    • Real-time Market Sentiment Analysis
    • News Integration and Impact Analysis
    • Social Sentiment Tracking
    • Institutional Trading Insights
    • Multi-source Data Integration
    """
    
    # Technology Stack Slide
    tech_stack_slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = tech_stack_slide.shapes.title
    content = tech_stack_slide.placeholders[1]
    title.text = "Technology Stack"
    content.text = """
    • Frontend: TypeScript, React
    • Backend: Node.js, TypeScript
    • APIs Integration:
        - Market Data APIs
        - News APIs
        - Trading APIs
    • Real-time Data Processing
    • Cloud Infrastructure
    """
    
    # Future Roadmap
    roadmap_slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = roadmap_slide.shapes.title
    content = roadmap_slide.placeholders[1]
    title.text = "Future Roadmap"
    content.text = """
    • AI-Powered Trading Signals
    • Machine Learning Risk Assessment
    • Extended Asset Classes
    • Mobile Application
    • Social Trading Features
    • Advanced Backtesting Engine
    """
    
    # Save the presentation
    prs.save('Algo360FX_Presentation.pptx')

if __name__ == "__main__":
    create_algo360fx_presentation()
